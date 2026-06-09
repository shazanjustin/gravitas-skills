"""
Generate one PPTX slide per post from a social media quarterly report Excel file.

Usage:
    python generate_slides.py --excel <path.xlsx> [--output <path.pptx>] [--sheet <name>] [--year <YYYY>]
"""
import argparse, io, re, sys, os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image, ImageDraw

# ── Colors ─────────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
ACCENT_RED = RGBColor(0xBB, 0x11, 0x11)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY   = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Column definitions ─────────────────────────────────────────────────────────
# (display_name, excel_col_idx)  None = not available in this platform → show '-'
IG_COLS = [
    ('Campaign name',       0),
    ('Placement',           1),
    ('Amount spent\n(MYR)', 4),
    ('Impressions',         5),
    ('Reach',               6),
    ('Frequency',           7),
    ('Post\ncomments',      9),
    ('Post\nshares',       10),
    ('Post\nreactions',    11),
    ('Post\nengagements',  12),
    ('Post\nsaves',        13),
    ('Clicks\n(all)',      14),
    ('Views',              15),
]

TT_COLS = [
    ('Campaign name',       0),
    ('Placement',           1),    # Advertising objective
    ('Amount spent\n(MYR)', 2),
    ('Impressions',         3),
    ('Reach',               4),
    ('Frequency',           5),
    ('Post\ncomments',      7),    # Paid comments
    ('Post\nshares',        8),    # Paid shares
    ('Post\nreactions',     6),    # Paid likes
    ('Post\nengagements', None),
    ('Post\nsaves',       None),
    ('Clicks\n(all)',       9),
    ('Views',              10),    # Video views
]

YT_COLS = [
    ('Campaign',            0),
    ('Objective',           1),    # Campaign type
    ('Impressions',         2),    # Impr.
    ('Reach',               3),
    ('Amount spent\n(MYR)', 4),
    ('Engagements',         5),
    ('TrueView\nviews',     6),
    ('Engagement\nrate',    7),
    ('Video played\nto 25%',  13),
    ('Video played\nto 50%',  14),
    ('Video played\nto 75%',  15),
    ('Video played\nto 100%', 16),
]

PLATFORM_META = {
    'instagram': ('INSTAGRAM PERFORMANCE', 'ig'),
    'tiktok':    ('TIKTOK PERFORMANCE',    'tt'),
    'youtube':   ('YOUTUBE PERFORMANCE',   'yt'),
}

# ── Value formatting ───────────────────────────────────────────────────────────
def fmt_val(val):
    if val is None:
        return '-'
    try:
        if pd.isna(val):
            return '-'
    except Exception:
        pass
    if isinstance(val, float):
        if val == 0:
            return '-'
        if val >= 1_000_000:
            return f'{val/1_000_000:.2f}M'
        if val >= 10_000:
            return f'{int(round(val)):,}'
        if val >= 100:
            return f'{val:,.2f}' if val != int(val) else f'{int(val):,}'
        if val < 1:
            return f'{val:.5f}'
        return f'{val:.2f}' if val != int(val) else str(int(val))
    return str(val) if str(val) != 'nan' else '-'


def cell_val(row, col_idx):
    if col_idx is None:
        return '-'
    try:
        return fmt_val(row.iloc[col_idx])
    except Exception:
        return '-'


def parse_period(campaign_name, year):
    m = re.search(r'\|(\d{1,2}-\w+)\|(\d{1,2}-\w+)$', campaign_name.strip())
    if m:
        return f'{m.group(1)} - {m.group(2)} {year}'
    return str(year)


def fk(n):
    if n >= 1_000_000: return f'{n/1_000_000:.1f}M'
    if n >= 1_000:     return f'{n/1_000:.1f}K'
    return str(int(round(n)))

# ── Cell styling ───────────────────────────────────────────────────────────────
def set_cell_border(cell, hex_color='CCCCCC', width_pt=0.5):
    emu = int(width_pt * 12700)
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn('a:tcBdr')):
        tcPr.remove(old)
    tcBdr = etree.SubElement(tcPr, qn('a:tcBdr'))
    for side in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln = etree.SubElement(tcBdr, qn(f'a:{side}'), w=str(emu))
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr'), val=hex_color)


def set_cell_fill(cell, rgb):
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb

# ── Platform logos (Pillow) ────────────────────────────────────────────────────
def make_ig_logo(size=96):
    grad = Image.new('RGBA', (size, size), (0, 0, 0, 0))
    gd = ImageDraw.Draw(grad)
    for y in range(size):
        r = int(193 + (252 - 193) * y / size)
        g = int(53  + (176 -  53) * y / size)
        b = int(132 + (69  - 132) * y / size)
        gd.rectangle([0, y, size - 1, y + 1], fill=(r, g, b, 255))
    mask = Image.new('L', (size, size), 0)
    ImageDraw.Draw(mask).rounded_rectangle([0, 0, size-1, size-1], radius=size//5, fill=255)
    grad.putalpha(mask)
    d = ImageDraw.Draw(grad)
    m, lw = size // 7, max(2, size // 20)
    d.rounded_rectangle([m, m, size-m-1, size-m-1], radius=size//10, outline='white', width=lw)
    cx, cy, lr = size//2, size//2 + size//18, size//5
    d.ellipse([cx-lr, cy-lr, cx+lr, cy+lr], outline='white', width=lw)
    fx, fy = size-m-size//7, m+size//10
    d.ellipse([fx-lw, fy-lw, fx+lw, fy+lw], fill='white')
    buf = io.BytesIO(); grad.save(buf, 'PNG'); return buf.getvalue()


def make_tt_logo(size=96):
    img = Image.new('RGBA', (size, size), (0, 0, 0, 255))
    d = ImageDraw.Draw(img)
    nx, ny = int(size*.35), int(size*.15)
    nw, nh, nr = int(size*.30), int(size*.55), int(size*.18)
    d.rectangle([nx+6, ny+4, nx+nw+6, ny+nh*3//4+4], fill=(105, 201, 208))
    d.ellipse([nx-nr+6, ny+nh-nr+4, nx+nr+6, ny+nh+nr+4], fill=(105, 201, 208))
    d.rectangle([nx-4, ny-4, nx+nw-4, ny+nh*3//4-4], fill=(255, 0, 80))
    d.ellipse([nx-nr-4, ny+nh-nr-4, nx+nr-4, ny+nh+nr-4], fill=(255, 0, 80))
    d.rectangle([nx, ny, nx+nw, ny+nh*3//4], fill='white')
    d.ellipse([nx-nr, ny+nh-nr, nx+nr, ny+nh+nr], fill='white')
    fx = nx + nw
    d.rectangle([fx, ny, fx+int(size*.2), ny+int(size*.25)], fill=(105, 201, 208))
    d.ellipse([fx, ny+int(size*.05), fx+int(size*.2), ny+int(size*.05)+int(size*.2)], fill=(105, 201, 208))
    buf = io.BytesIO(); img.save(buf, 'PNG'); return buf.getvalue()


def make_yt_logo(size=96):
    h = int(size * 0.7)
    img = Image.new('RGBA', (size, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle([0, 0, size-1, h-1], radius=h//6, fill=(255, 0, 0, 255))
    tx, ty = int(size*.38), int(h*.25)
    d.polygon([(tx, ty), (tx, ty+int(h*.50)), (tx+int(size*.30), ty+int(h*.25))], fill='white')
    buf = io.BytesIO(); img.save(buf, 'PNG'); return buf.getvalue()

# ── Excel parser ───────────────────────────────────────────────────────────────
STOP_WORDS = {'Campaign name', 'Campaign', 'JANUARY', 'FEBRUARY', 'MARCH', 'APRIL',
              'MAY', 'JUNE', 'JULY', 'AUGUST', 'SEPTEMBER', 'OCTOBER', 'NOVEMBER',
              'DECEMBER', 'BY MONTHLY', 'TOTAL', 'META', 'TIKTOK', 'YOUTUBE',
              'OVERALL CAMPAIGN'}


def detect_sheet(excel_path):
    xl = pd.ExcelFile(excel_path)
    return xl.sheet_names[0]


def parse_excel(excel_path, sheet_name):
    df = pd.read_excel(excel_path, sheet_name=sheet_name, header=None)
    blocks = []
    month, i, n = None, 0, len(df)

    while i < n:
        row = df.iloc[i]
        v0 = str(row.iloc[0]).strip() if pd.notna(row.iloc[0]) else ''
        v1 = str(row.iloc[1]).strip() if pd.notna(row.iloc[1]) else ''

        if v0.upper() in {'JANUARY','FEBRUARY','MARCH','APRIL','MAY','JUNE',
                          'JULY','AUGUST','SEPTEMBER','OCTOBER','NOVEMBER','DECEMBER'}:
            month = v0.strip()
            i += 1
            continue

        if v0 in ('Campaign name', 'Campaign') and \
                v1 in ('Placement', 'Advertising objective', 'Campaign type'):
            if v1 == 'Placement':
                platform, cols = 'instagram', IG_COLS
            elif v1 == 'Advertising objective':
                platform, cols = 'tiktok', TT_COLS
            else:
                platform, cols = 'youtube', YT_COLS

            j = i + 1
            if platform == 'instagram':
                data_rows, cname = [], None
                while j < n:
                    dr = df.iloc[j]
                    dv0 = str(dr.iloc[0]).strip() if pd.notna(dr.iloc[0]) else ''
                    dv1 = str(dr.iloc[1]).strip() if pd.notna(dr.iloc[1]) else ''
                    if not dv0 and not dv1:
                        break
                    if dv0 in STOP_WORDS or dv0.startswith('Link:') or dv0.startswith('*'):
                        break
                    if dv0:
                        cname = dv0
                    data_rows.append(dr)
                    j += 1
                if data_rows and cname:
                    blocks.append({'month': month, 'platform': platform,
                                   'campaign_name': cname, 'data_rows': data_rows, 'cols': cols})
            else:
                while j < n:
                    dr = df.iloc[j]
                    dv0 = str(dr.iloc[0]).strip() if pd.notna(dr.iloc[0]) else ''
                    dv1 = str(dr.iloc[1]).strip() if pd.notna(dr.iloc[1]) else ''
                    if not dv0 and not dv1:
                        break
                    if dv0 in STOP_WORDS or dv0.startswith('Link:') or dv0.startswith('*'):
                        break
                    if dv0:
                        blocks.append({'month': month, 'platform': platform,
                                       'campaign_name': dv0, 'data_rows': [dr], 'cols': cols})
                    j += 1
            i = j
        else:
            i += 1

    return blocks

# ── Insights ───────────────────────────────────────────────────────────────────
def gen_insights(block):
    rows, platform = block['data_rows'], block['platform']

    def safe(row, idx):
        try:
            v = row.iloc[idx]
            return float(v) if pd.notna(v) else 0.0
        except Exception:
            return 0.0

    bullets = []
    if platform == 'instagram':
        spend = sum(safe(r,4) for r in rows)
        impr  = sum(safe(r,5) for r in rows)
        reach = sum(safe(r,6) for r in rows)
        views = sum(safe(r,15) for r in rows)
        bullets.append(
            f'With a total spend of RM{spend:,.2f}, the campaign delivered '
            f'{fk(impr)} impressions, {fk(reach)} reach and {fk(views)} views.'
        )
        if len(rows) > 1:
            ranked = sorted(rows, key=lambda r: safe(r,6), reverse=True)
            p1 = str(ranked[0].iloc[1]) if pd.notna(ranked[0].iloc[1]) else 'Top placement'
            bullets.append(
                f'{p1} was the strongest placement, generating the highest reach '
                f'({fk(safe(ranked[0],6))}), engagements ({fk(safe(ranked[0],12))}) '
                f'and clicks ({fk(safe(ranked[0],14))}).'
            )
            if len(rows) >= 3:
                p2 = str(ranked[1].iloc[1]) if pd.notna(ranked[1].iloc[1]) else 'Second placement'
                p3 = str(ranked[2].iloc[1]) if pd.notna(ranked[2].iloc[1]) else 'third placement'
                bullets.append(
                    f'{p2} performed as the second strongest placement, while '
                    f'{p3} added supporting reach at lower spend.'
                )
            elif len(rows) == 2:
                p2 = str(ranked[1].iloc[1]) if pd.notna(ranked[1].iloc[1]) else 'Second placement'
                bullets.append(f'{p2} supported overall reach with lower spend.')

    elif platform == 'tiktok':
        r = rows[0]
        bullets.append(
            f'With a spend of RM{safe(r,2):,.2f}, the TikTok post delivered '
            f'{fk(safe(r,3))} impressions and {fk(safe(r,10))} video views.'
        )
        bullets.append(f'Reached {fk(safe(r,4))} unique users with a frequency of {safe(r,5):.2f}.')
        if safe(r,9):
            bullets.append(f'Generated {fk(safe(r,9))} total clicks across the campaign period.')

    elif platform == 'youtube':
        r = rows[0]
        bullets.append(
            f'With a spend of RM{safe(r,4):,.2f}, the YouTube campaign delivered '
            f'{fk(safe(r,2))} impressions and {fk(safe(r,6))} TrueView views.'
        )
        if safe(r,5):
            er = safe(r,7)
            bullets.append(
                f'Total engagements reached {fk(safe(r,5))}'
                + (f' with an engagement rate of {er:.2%}.' if er else '.')
            )
        p25, p100 = safe(r,13), safe(r,16)
        if p25:
            bullets.append(
                f'Video completion funnel: 25% at {p25:.1%}'
                + (f', 100% at {p100:.1%}.' if p100 else '.')
            )

    return bullets[:3]

# ── Slide builder ──────────────────────────────────────────────────────────────
def build_slide(prs, block, logos, year):
    platform  = block['platform']
    cname     = block['campaign_name']
    data_rows = block['data_rows']
    cols      = block['cols']
    title_text, logo_key = PLATFORM_META[platform]

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    hdr_h = Inches(0.72)
    hdr = slide.shapes.add_shape(1, 0, 0, SLIDE_W, hdr_h)
    hdr.fill.solid(); hdr.fill.fore_color.rgb = BLACK
    hdr.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.2), Inches(0.04), Inches(10.5), hdr_h - Inches(0.08))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = title_text
    run.font.bold = True; run.font.size = Pt(26)
    run.font.color.rgb = WHITE; run.font.name = 'Arial'

    # Platform logo
    logo_bytes = logos.get(logo_key)
    if logo_bytes:
        lw = lh = Inches(0.65)
        slide.shapes.add_picture(
            io.BytesIO(logo_bytes),
            SLIDE_W - lw - Inches(0.15), Inches(0.04), lw, lh
        )

    # Reporting period
    period = parse_period(cname, year)
    st = slide.shapes.add_textbox(Inches(0.2), Inches(0.78), Inches(9), Inches(0.38))
    tf = st.text_frame; p = tf.paragraphs[0]
    run = p.add_run(); run.text = f'Reporting Period : {period}'
    run.font.size = Pt(11); run.font.italic = True
    run.font.color.rgb = ACCENT_RED; run.font.name = 'Arial'

    # Table
    n_data = len(data_rows)
    n_cols = len(cols)
    header_rh = Inches(0.52)
    data_rh   = min(Inches(0.54), Inches(3.6 / max(n_data, 1)))
    table_h   = header_rh + data_rh * n_data
    table_top = Inches(1.22)
    usable_w  = SLIDE_W - Inches(0.3)

    widths_pct = []
    for name, _ in cols:
        if 'Campaign' in name:                             widths_pct.append(20)
        elif 'Placement' in name or 'Objective' in name:  widths_pct.append(8)
        elif 'Engagement rate' in name:                    widths_pct.append(6)
        else:                                              widths_pct.append(100 / max(n_cols-2, 1) * 0.72)
    total_pct  = sum(widths_pct)
    col_widths = [int(usable_w * p / total_pct) for p in widths_pct]
    col_widths[-1] += int(usable_w) - sum(col_widths)

    tbl_shape = slide.shapes.add_table(
        n_data + 1, n_cols, Inches(0.15), table_top, sum(col_widths), table_h
    )
    tbl = tbl_shape.table
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = w
    tbl.rows[0].height = header_rh
    for ri in range(1, n_data + 1):
        tbl.rows[ri].height = data_rh

    for ci, (col_name, _) in enumerate(cols):
        cell = tbl.cell(0, ci)
        set_cell_fill(cell, ACCENT_RED)
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = col_name.replace('\n', ' ')
        run.font.bold = True; run.font.size = Pt(7)
        run.font.color.rgb = WHITE; run.font.name = 'Arial'
        set_cell_border(cell, 'FFFFFF', 0.75)

    carried_name = None
    for ri, dr in enumerate(data_rows):
        dv0 = str(dr.iloc[0]).strip() if pd.notna(dr.iloc[0]) else ''
        if dv0 and dv0 != 'nan':
            carried_name = dv0
        for ci, (col_name, xcol) in enumerate(cols):
            cell = tbl.cell(ri + 1, ci)
            set_cell_fill(cell, LIGHT_GRAY if ri % 2 == 1 else WHITE)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            if ci == 0:
                run.text = carried_name or ''
                run.font.size = Pt(6.5)
            else:
                run.text = cell_val(dr, xcol)
                run.font.size = Pt(7.5)
            run.font.name = 'Arial'
            run.font.color.rgb = DARK_GRAY
            set_cell_border(cell, 'CCCCCC', 0.5)

    # Insights box
    ins_top = table_top + table_h + Inches(0.15)
    ins_h   = SLIDE_H - ins_top - Inches(0.15)
    ins_w   = Inches(7.8)
    ins_bg  = slide.shapes.add_shape(1, Inches(0.15), ins_top, ins_w, ins_h)
    ins_bg.fill.solid(); ins_bg.fill.fore_color.rgb = WHITE
    ins_bg.line.color.rgb = ACCENT_RED; ins_bg.line.width = Pt(1.5)
    ins_tb = slide.shapes.add_textbox(
        Inches(0.3), ins_top + Inches(0.1), ins_w - Inches(0.3), ins_h - Inches(0.2)
    )
    tf = ins_tb.text_frame; tf.word_wrap = True
    for idx, bullet in enumerate(gen_insights(block)):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        if idx > 0:
            p.space_before = Pt(5)
        run = p.add_run(); run.text = f'•  {bullet}'
        run.font.size = Pt(8.5); run.font.name = 'Arial'
        run.font.color.rgb = DARK_GRAY

    # Thumbnail placeholder
    thumb_l = Inches(0.15) + ins_w + Inches(0.1)
    thumb_w = SLIDE_W - thumb_l - Inches(0.15)
    th_bg = slide.shapes.add_shape(1, thumb_l, ins_top, thumb_w, ins_h)
    th_bg.fill.solid(); th_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    th_bg.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    th_tb = slide.shapes.add_textbox(
        thumb_l, ins_top + ins_h / 2 - Inches(0.15), thumb_w, Inches(0.3)
    )
    tf = th_tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = '[Thumbnail]'
    run.font.size = Pt(9); run.font.italic = True
    run.font.color.rgb = MID_GRAY; run.font.name = 'Arial'

# ── Entry point ────────────────────────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser(description='Generate PPTX slides from social report Excel')
    ap.add_argument('--excel',  required=True, help='Path to the .xlsx report file')
    ap.add_argument('--output', default=None,  help='Output .pptx path (default: alongside Excel)')
    ap.add_argument('--sheet',  default=None,  help='Sheet name (default: auto-detect first sheet)')
    ap.add_argument('--year',   default=None,  help='Report year e.g. 2026 (default: auto from filename)')
    args = ap.parse_args()

    excel_path = args.excel
    sheet_name = args.sheet or detect_sheet(excel_path)

    if args.year:
        year = args.year
    else:
        m = re.search(r'20\d\d', os.path.basename(excel_path))
        year = m.group(0) if m else '2026'

    if args.output:
        output_path = args.output
    else:
        base = os.path.splitext(excel_path)[0]
        output_path = base + '_slides.pptx'

    print(f'Sheet   : {sheet_name}')
    print(f'Year    : {year}')
    print(f'Output  : {output_path}')
    print()

    blocks = parse_excel(excel_path, sheet_name)
    print(f'Found {len(blocks)} posts:')
    for b in blocks:
        print(f"  [{b['month'] or '?':>9}] {b['platform'].upper():<10} {b['campaign_name'][:65]}...")

    logos = {'ig': make_ig_logo(), 'tt': make_tt_logo(), 'yt': make_yt_logo()}

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for idx, block in enumerate(blocks, 1):
        print(f'  Slide {idx:>2}/{len(blocks)}  [{block["month"]}] {block["platform"].upper()}')
        build_slide(prs, block, logos, year)

    prs.save(output_path)
    print(f'\nSaved: {output_path}')


if __name__ == '__main__':
    main()
